"""
ticket_dispatcher.py — Dispatch de acciones para agentes y herramientas.
"""

import asyncio
import base64
import contextlib
import hashlib
import io
import json
import mimetypes
import multiprocessing
import os
import queue
import resource
import signal
import subprocess
import time
import traceback
import uuid
from difflib import SequenceMatcher
from typing import Any

from dotenv import load_dotenv
from openai import AsyncOpenAI

from tools.memoryTools.RAG_memory import MemoryRag
from tools.memoryTools.semantic_splitter import count_tokens
# ── Env / Clients / RAG ─────────────────────────────────────────────

load_dotenv()
RAG = MemoryRag()
openai_client = AsyncOpenAI(api_key=os.environ["OPENAI_API_KEY"])
user_preferences_path = os.getenv("USER_PREFERENCES_PATH", "/tmp/user_preferences.txt")

# ── Constantes ───────────────────────────────────────────────────────

SKIP_DIRS = {".git", "__pycache__", ".venv", "venv", "node_modules"}
TEXT_EXTENSIONS = {".txt", ".md", ".py", ".json", ".csv", ".log", ".yaml", ".yml", ".sql"}
DEFAULT_MAX_CHARS = 200_000
DEFAULT_SEARCH_LIMIT = 50
DEFAULT_PLAYWRIGHT_MAX_CHARS = 12_000
DEFAULT_PLAYWRIGHT_ELEMENT_LIMIT = 40
DEFAULT_PLAYWRIGHT_SESSION_TIMEOUT = 300
MAX_PLAYWRIGHT_SESSION_TIMEOUT = 2_700
FILE_HASH_ALGORITHM = "md5"
HASH_CHUNK_BYTES = 1024 * 1024

ALLOWED_WRITE_ROOTS = [
    os.path.expanduser("~/Downloads"),
    os.path.expanduser("~/Documents"),
    os.path.expanduser("~/Desktop"),
    os.path.expanduser("~"),
    "/mnt/d",
    "/mnt/c",
    "/tmp",
]

BLOCKED_COMMANDS = {
    "rm -rf /", "rm -rf /*", "mkfs", "dd if=", ":(){", "fork",
    "chmod -R 777 /", "chown -R", "shutdown", "reboot", "halt",
    "poweroff", "init 0", "init 6", "killall", "kill -9 1",
    "> /dev/sda", "mv / ", "wget | sh", "curl | sh", "curl | bash",
    "wget | bash", "python -c", "python3 -c",
}

BLOCKED_PREFIXES = {"sudo", "su", "doas"}

SAFE_BUILTINS = {
    "abs": abs, "min": min, "max": max, "sum": sum,
    "len": len, "range": range, "print": print,
    "int": int, "float": float, "str": str, "bool": bool,
    "list": list, "dict": dict, "tuple": tuple, "set": set,
    "sorted": sorted, "reversed": reversed, "enumerate": enumerate,
    "zip": zip, "map": map, "filter": filter,
    "isinstance": isinstance, "type": type,
    "round": round, "pow": pow, "divmod": divmod,
    "True": True, "False": False, "None": None,
}


# ── Utilidad de rutas ────────────────────────────────────────────────

def _resolve_path(path: str) -> str:
    """Expande ~ y variables de entorno, luego resuelve a ruta absoluta."""
    return os.path.realpath(os.path.expandvars(os.path.expanduser(path)))


def _file_hash(path: str) -> dict:
    """Calcula hash en streaming para no cargar archivos grandes en memoria."""
    hasher = hashlib.md5()

    with open(path, "rb") as f:
        for chunk in iter(lambda: f.read(HASH_CHUNK_BYTES), b""):
            hasher.update(chunk)

    stat = os.stat(path)
    return {
        "algorithm": FILE_HASH_ALGORITHM,
        "value": hasher.hexdigest(),
        "bytes": stat.st_size,
        "mtime_ns": stat.st_mtime_ns,
    }


# ── Herramientas: Archivos ──────────────────────────────────────────

def action_search_files(body: dict) -> dict:
    """Busca archivos por nombre. Requiere: query. Opcional: root, limit."""
    query = (body.get("query") or "").strip()
    if not query:
        return {"success": False, "error": "missing required field: query", "results": []}

    root = body.get("root") or "."
    limit = max(int(body.get("limit") or DEFAULT_SEARCH_LIMIT), 1)
    terms = query.split()
    q_cf = query.casefold()

    def score(name: str) -> float:
        n_cf = name.casefold()
        exact = 1.0 if q_cf in n_cf else 0.0
        coverage = sum(t.casefold() in n_cf for t in terms) / max(len(terms), 1)
        fuzzy = SequenceMatcher(a=q_cf, b=n_cf).ratio()
        return 0.9 * fuzzy + 0.8 * coverage + 0.6 * exact

    results = []
    try:
        for dirpath, dirnames, filenames in os.walk(root, followlinks=False):
            dirnames[:] = [d for d in dirnames if d not in SKIP_DIRS]
            for fn in filenames:
                s = score(fn)
                if s > 0.25:
                    results.append({"path": os.path.join(dirpath, fn), "filename": fn, "score": s})

        results.sort(key=lambda r: r["score"], reverse=True)
        return {
            "success": True,
            "query": query,
            "root": root,
            "limit": limit,
            "total_matches": len(results),
            "remaining_results": max(len(results) - limit, 0),
            "results": results[:limit],
        }
    except Exception as e:
        return {"success": False, "error": str(e), "results": []}


def action_read_file(body: dict) -> dict:
    """Lee un archivo. Requiere: path. Opcional: max_chars."""
    path = body.get("path")
    if not path:
        return {"success": False, "error": "missing required field: path"}

    path = _resolve_path(path)
    if not os.path.exists(path):
        return {"success": False, "error": f"file not found: {path}"}

    max_chars = max(int(body.get("max_chars") or DEFAULT_MAX_CHARS), 1)
    ext = os.path.splitext(path)[1].lower()
    mime = mimetypes.guess_type(path)[0] or "application/octet-stream"

    try:
        reader = _FILE_READERS.get(ext)
        if reader:
            result = reader(path, mime, max_chars)
        else:
            result = _read_binary(path, mime)
        if result.get("success"):
            result["file_hash"] = _file_hash(path)
        return result
    except Exception as e:
        return {"success": False, "path": path, "mime": mime, "error": str(e)}


def action_file_hash(body: dict) -> dict:
    """Devuelve una huella eficiente del archivo. Requiere: path."""
    path = body.get("path")
    if not path:
        return {"success": False, "error": "missing required field: path"}

    real_path = _resolve_path(path)
    if not os.path.isfile(real_path):
        return {"success": False, "error": f"file not found: {real_path}", "path": real_path}

    compare_to = (body.get("compare_to") or "").strip() or None

    try:
        file_hash = _file_hash(real_path)
    except Exception as e:
        return {"success": False, "path": real_path, "error": str(e)}

    result = {
        "success": True,
        "path": real_path,
        "file_hash": file_hash,
    }
    if compare_to is not None:
        result["changed"] = file_hash["value"] != compare_to
    return result


def action_write_file(body: dict) -> dict:
    """
    Escribe contenido a un archivo de texto.
    Requiere: path (str), content (str).
    Opcional: mode ('w'|'a'), encoding (default 'utf-8').
    """
    path = body.get("path")
    content = body.get("content")
    if not path:
        return {"success": False, "error": "missing required field: path"}
    if content is None:
        return {"success": False, "error": "missing required field: content"}

    mode = body.get("mode", "w")
    if mode not in ("w", "a"):
        return {"success": False, "error": "mode must be 'w' or 'a'"}
    encoding = body.get("encoding", "utf-8")

    # ✅ FIX: expandir ~ y $HOME antes de validar
    real_path = _resolve_path(path)

    if not any(real_path.startswith(os.path.realpath(root)) for root in ALLOWED_WRITE_ROOTS):
        return {
            "success": False,
            "error": f"write not allowed outside: {ALLOWED_WRITE_ROOTS}",
            "path": real_path,
        }

    try:
        os.makedirs(os.path.dirname(real_path), exist_ok=True)
        with open(real_path, mode, encoding=encoding) as f:
            f.write(content)
        return {
            "success": True,
            "path": real_path,
            "bytes_written": len(content.encode(encoding)),
            "file_hash": _file_hash(real_path),
        }
    except Exception as e:
        return {"success": False, "path": real_path, "error": str(e)}


def action_edit_file(body: dict) -> dict:
    """
    Modifica un archivo reemplazando old_text por new_text (search-and-replace).
    Requiere: path, old_text, new_text.
    Opcional: replace_all (bool, default False).
    """
    path = body.get("path")
    old_text = body.get("old_text")
    new_text = body.get("new_text")

    if not path:
        return {"success": False, "error": "missing required field: path"}
    if old_text is None:
        return {"success": False, "error": "missing required field: old_text"}
    if new_text is None:
        return {"success": False, "error": "missing required field: new_text"}

    real_path = _resolve_path(path)

    if not any(real_path.startswith(os.path.realpath(root)) for root in ALLOWED_WRITE_ROOTS):
        return {
            "success": False,
            "error": f"edit not allowed outside: {ALLOWED_WRITE_ROOTS}",
            "path": real_path,
        }

    if not os.path.isfile(real_path):
        return {"success": False, "error": f"file not found: {real_path}"}

    replace_all = bool(body.get("replace_all", False))

    try:
        with open(real_path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        count = content.count(old_text)

        if count == 0:
            return {
                "success": False,
                "error": "old_text not found in file",
                "path": real_path,
            }

        if not replace_all and count > 1:
            return {
                "success": False,
                "error": f"old_text found {count} times — set replace_all=true or provide more context to make it unique",
                "path": real_path,
                "occurrences": count,
            }

        new_content = content.replace(old_text, new_text) if replace_all else content.replace(old_text, new_text, 1)

        with open(real_path, "w", encoding="utf-8") as f:
            f.write(new_content)

        replacements = count if replace_all else 1
        return {
            "success": True,
            "path": real_path,
            "replacements": replacements,
            "file_hash": _file_hash(real_path),
        }
    except Exception as e:
        return {"success": False, "path": real_path, "error": str(e)}


# ── Lectores por tipo ────────────────────────────────────────────────

def _clip_text_content(content: str, max_chars: int) -> tuple[str, bool, int]:
    shown = content[:max_chars]
    remaining_chars = max(len(content) - len(shown), 0)
    return shown, remaining_chars > 0, remaining_chars


def _count_remaining_tokens(content: str, max_chars: int) -> int:
    if len(content) <= max_chars:
        return 0
    return count_tokens(content[max_chars:])


def _read_text(path: str, mime: str, max_chars: int) -> dict:
    with open(path, "r", encoding="utf-8", errors="replace") as f:
        content = f.read()
    clipped, truncated, remaining_chars = _clip_text_content(content, max_chars)
    return {
        "success": True, "path": path, "mime": mime, "kind": "text",
        "truncated": truncated, "remaining_chars": remaining_chars,
        "remaining_tokens": _count_remaining_tokens(content, max_chars),
        "content": clipped,
    }


def _read_pdf(path: str, mime: str, max_chars: int) -> dict:
    from PyPDF2 import PdfReader
    reader = PdfReader(path)
    content = []
    for page in reader.pages:
        txt = page.extract_text() or ""
        if not txt:
            continue
        content.append(txt)
    full_content = "\n".join(content)
    clipped, truncated, remaining_chars = _clip_text_content(full_content, max_chars)
    return {
        "success": True, "path": path, "mime": mime, "kind": "pdf_text",
        "pages": len(reader.pages), "truncated": truncated,
        "remaining_chars": remaining_chars,
        "remaining_tokens": _count_remaining_tokens(full_content, max_chars),
        "content": clipped,
    }


def _read_docx(path: str, mime: str, max_chars: int) -> dict:
    from docx import Document
    doc = Document(path)
    parts = []
    for p in doc.paragraphs:
        t = (p.text or "").strip()
        if not t:
            continue
        parts.append(t + "\n")
    full_content = "".join(parts)
    clipped, truncated, remaining_chars = _clip_text_content(full_content, max_chars)
    return {
        "success": True, "path": path, "mime": mime, "kind": "docx_text",
        "truncated": truncated, "remaining_chars": remaining_chars,
        "remaining_tokens": _count_remaining_tokens(full_content, max_chars),
        "content": clipped,
    }


def _read_xlsx(path: str, mime: str, max_chars: int) -> dict:
    import openpyxl
    wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
    sheets, total = {}, 0
    omitted_parts = []
    for ws in wb.worksheets:
        rows = []
        for row in ws.iter_rows(values_only=True):
            line = "\t".join("" if v is None else str(v) for v in row).rstrip()
            if not line:
                continue
            if total < max_chars:
                remaining_budget = max_chars - total
                rows.append(line[:remaining_budget])
                omitted = line[remaining_budget:]
                if omitted:
                    omitted_parts.append(omitted)
            else:
                omitted_parts.append(line)
            total += len(line) + 1
        if rows:
            sheets[ws.title] = rows
    return {
        "success": True, "path": path, "mime": mime, "kind": "xlsx_text",
        "sheets": list(sheets.keys()), "truncated": total > max_chars,
        "remaining_chars": max(total - max_chars, 0),
        "remaining_tokens": count_tokens("\n".join(omitted_parts)) if omitted_parts else 0,
        "content": sheets,
    }


def _read_pptx(path: str, mime: str, max_chars: int) -> dict:
    from pptx import Presentation
    prs = Presentation(path)
    texts = []
    for si, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            t = getattr(shape, "text", "").strip()
            if not t:
                continue
            block = f"[slide {si}] {t}\n"
            texts.append(block)
    full_content = "".join(texts)
    clipped, truncated, remaining_chars = _clip_text_content(full_content, max_chars)
    return {
        "success": True, "path": path, "mime": mime, "kind": "pptx_text",
        "slides": len(prs.slides), "truncated": truncated,
        "remaining_chars": remaining_chars,
        "remaining_tokens": _count_remaining_tokens(full_content, max_chars),
        "content": clipped,
    }


def _read_binary(path: str, mime: str) -> dict:
    size = os.path.getsize(path)
    return {
        "success": True, "path": path, "mime": mime, "kind": "binary",
        "bytes": size,
        "preview_omitted": True,
        "omission_reason": "binary preview disabled to avoid large base64 payloads in tool outputs",
    }


_FILE_READERS = {
    **{ext: _read_text for ext in TEXT_EXTENSIONS},
    ".pdf": _read_pdf,
    ".docx": _read_docx,
    ".xlsx": _read_xlsx,
    ".pptx": _read_pptx,
}


# ── Herramienta: Ejecutar comandos de terminal ──────────────────────

def _is_command_blocked(command: str) -> str | None:
    cmd_lower = command.strip().lower()
    first_word = cmd_lower.split()[0] if cmd_lower else ""
    for prefix in BLOCKED_PREFIXES:
        if first_word == prefix.strip():
            return f"comando bloqueado: no se permite '{prefix.strip()}'"
    for pattern in BLOCKED_COMMANDS:
        if pattern in cmd_lower:
            return f"comando bloqueado: patrón peligroso detectado '{pattern}'"
    return None


async def action_run_command(body: dict) -> dict:
    """
    Ejecuta un comando de terminal (bash).
    Requiere: command (str).
    Opcional: timeout (int, default 15, max 30), workdir (str).
    """
    command = (body.get("command") or "").strip()
    if not command:
        return {"success": False, "error": "missing required field: command"}

    timeout = min(int(body.get("timeout", 15)), 30)

    # ✅ FIX: expandir ~ y $HOME en workdir
    workdir = _resolve_path(body.get("workdir") or "~")

    blocked = _is_command_blocked(command)
    if blocked:
        return {"success": False, "error": blocked, "command": command}

    try:
        result = subprocess.run(
            command,
            shell=True,
            executable="/bin/bash",
            capture_output=True,
            text=True,
            timeout=timeout,
            cwd=workdir,
            env={**os.environ, "LC_ALL": "C.UTF-8"},
        )

        stdout = result.stdout[-DEFAULT_MAX_CHARS:] if len(result.stdout) > DEFAULT_MAX_CHARS else result.stdout
        stderr = result.stderr[-50_000:] if len(result.stderr) > 50_000 else result.stderr

        return {
            "success": result.returncode == 0,
            "command": command,
            "returncode": result.returncode,
            "stdout": stdout,
            "stderr": stderr,
        }

    except subprocess.TimeoutExpired:
        return {"success": False, "error": "timeout exceeded", "command": command, "timeout": timeout}
    except Exception as e:
        return {"success": False, "error": str(e), "command": command}


# ── Herramienta: Sandbox Python ─────────────────────────────────────

def _apply_sandbox_limits(memory_mb=512, cpu_seconds=10, max_files=64):
    mem = memory_mb * 1024 * 1024
    resource.setrlimit(resource.RLIMIT_AS, (mem, mem))
    resource.setrlimit(resource.RLIMIT_CPU, (cpu_seconds, cpu_seconds))
    resource.setrlimit(resource.RLIMIT_NOFILE, (max_files, max_files))
    resource.setrlimit(resource.RLIMIT_CORE, (0, 0))
    os.umask(0o077)
    try:
        import ctypes
        ctypes.CDLL("libc.so.6").prctl(1, signal.SIGKILL)
    except Exception:
        pass


def _run_sandboxed(code: str, result_queue: multiprocessing.Queue):
    try:
        _apply_sandbox_limits()
        safe_globals = {"__builtins__": SAFE_BUILTINS}
        buf = io.StringIO()
        with contextlib.redirect_stdout(buf):
            exec(code, safe_globals)

        result = safe_globals.get("RESULT", buf.getvalue().strip() or None)
        if result is not None:
            try:
                json.dumps(result)
            except TypeError:
                result = repr(result)

        result_queue.put((True, buf.getvalue(), result, None))
    except Exception as e:
        result_queue.put((False, "", None, traceback.format_exc()))


def action_run_python(body: dict) -> dict:
    """Ejecuta código Python en sandbox. Requiere: code. Opcional: timeout (max 10s)."""
    code = body.get("code", "")
    timeout = min(int(body.get("timeout", 5)), 10)

    q = multiprocessing.Queue()
    p = multiprocessing.Process(target=_run_sandboxed, args=(code, q))

    start = time.perf_counter()
    p.start()
    p.join(timeout)
    elapsed = time.perf_counter() - start

    if p.is_alive():
        p.kill()
        p.join()
        return {"success": False, "error": "timeout exceeded", "execution_time": elapsed}

    try:
        success, stdout, result, tb = q.get_nowait()
    except queue.Empty:
        return {"success": False, "error": "no result returned", "execution_time": elapsed}

    resp = {"success": success, "stdout": stdout, "execution_time": elapsed}
    if success:
        resp["result"] = result
    else:
        resp["error"] = tb
    return resp


async def action_web_fetch(body: dict) -> dict:
    """
    Descarga el contenido de una URL vía HTTP.
    Requiere: url (str).
    Opcional: method ('GET'|'POST', default 'GET'), headers (dict),
              data (dict|str, para POST), timeout (int, default 15),
              max_chars (int).
    """
    import httpx

    url = (body.get("url") or "").strip()
    if not url:
        return {"success": False, "error": "missing required field: url"}

    method = (body.get("method") or "GET").upper()
    if method not in ("GET", "POST"):
        return {"success": False, "error": "method must be 'GET' or 'POST'"}

    headers = body.get("headers") or {}
    data = body.get("data")
    timeout = min(int(body.get("timeout") or 15), 60)
    max_chars = max(int(body.get("max_chars") or DEFAULT_MAX_CHARS), 1)

    try:
        async with httpx.AsyncClient(follow_redirects=True, timeout=timeout) as client:
            if method == "POST":
                if isinstance(data, dict):
                    response = await client.post(url, json=data, headers=headers)
                else:
                    response = await client.post(url, content=data or "", headers=headers)
            else:
                response = await client.get(url, headers=headers)

        content_type = response.headers.get("content-type", "")
        text = response.text
        content, truncated, remaining_chars = _clip_text_content(text, max_chars)

        return {
            "success": response.is_success,
            "url": str(response.url),
            "status_code": response.status_code,
            "content_type": content_type,
            "truncated": truncated,
            "remaining_chars": remaining_chars,
            "remaining_tokens": _count_remaining_tokens(text, max_chars),
            "content": content,
        }
    except httpx.TimeoutException:
        return {"success": False, "error": "timeout exceeded", "url": url}
    except Exception as e:
        return {"success": False, "error": str(e), "url": url}


def _coerce_bool(value, default: bool = False) -> bool:
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    if isinstance(value, str):
        return value.strip().lower() in {"1", "true", "yes", "y", "si", "sí"}
    return bool(value)


def _truncate_text(text: str, max_chars: int) -> tuple[str, bool, int]:
    return _clip_text_content(text, max_chars)


def _playwright_output_dir(body: dict) -> str:
    requested = (body.get("output_dir") or "").strip()
    base = requested or os.getenv("PLAYWRIGHT_OUTPUT_DIR") or os.path.join("/tmp", "planner_playwright")
    output_dir = _resolve_path(base)
    os.makedirs(output_dir, exist_ok=True)
    return output_dir


async def _collect_playwright_snapshot(
    page,
    max_chars: int,
    include_text: bool,
    include_elements: bool,
    include_html: bool,
) -> dict:
    snapshot: dict = {
        "title": await page.title(),
        "url": page.url,
    }

    if include_text:
        text = await page.locator("body").inner_text(timeout=3000)
        (
            snapshot["text"],
            snapshot["text_truncated"],
            snapshot["text_remaining_chars"],
        ) = _truncate_text(text, max_chars)
        snapshot["text_remaining_tokens"] = _count_remaining_tokens(text, max_chars)

    if include_elements:
        elements = await page.evaluate(
            """(limit) => {
                const visibleText = (el) => (el.innerText || el.value || el.getAttribute('aria-label') || el.getAttribute('title') || '').trim();
                const clip = (s, n = 160) => s && s.length > n ? s.slice(0, n) : s;
                const links = Array.from(document.querySelectorAll('a[href]')).slice(0, limit).map((a) => ({
                    text: clip(visibleText(a)),
                    href: a.href
                })).filter((x) => x.text || x.href);
                const buttons = Array.from(document.querySelectorAll('button, input[type="button"], input[type="submit"], [role="button"]')).slice(0, limit).map((b) => ({
                    text: clip(visibleText(b)),
                    type: b.getAttribute('type') || b.tagName.toLowerCase()
                })).filter((x) => x.text || x.type);
                const fields = Array.from(document.querySelectorAll('input, textarea, select')).slice(0, limit).map((el) => ({
                    name: el.getAttribute('name') || '',
                    id: el.id || '',
                    type: el.getAttribute('type') || el.tagName.toLowerCase(),
                    placeholder: el.getAttribute('placeholder') || '',
                    label: clip((el.labels && el.labels[0] && el.labels[0].innerText) || el.getAttribute('aria-label') || '')
                }));
                const forms = Array.from(document.querySelectorAll('form')).slice(0, limit).map((form) => ({
                    action: form.action || '',
                    method: form.method || 'get',
                    fields: Array.from(form.querySelectorAll('input, textarea, select')).slice(0, 20).map((el) => el.getAttribute('name') || el.id || el.getAttribute('placeholder') || el.tagName.toLowerCase()).filter(Boolean)
                }));
                return {links, buttons, fields, forms};
            }""",
            DEFAULT_PLAYWRIGHT_ELEMENT_LIMIT,
        )
        snapshot["elements"] = elements

    if include_html:
        html = await page.content()
        (
            snapshot["html"],
            snapshot["html_truncated"],
            snapshot["html_remaining_chars"],
        ) = _truncate_text(html, max_chars)
        snapshot["html_remaining_tokens"] = _count_remaining_tokens(html, max_chars)

    return snapshot


async def _save_playwright_screenshot(page, body: dict, full_page: bool) -> dict:
    output_dir = _playwright_output_dir(body)
    filename = f"playwright_{int(time.time())}_{uuid.uuid4().hex[:8]}.png"
    path = os.path.join(output_dir, filename)
    screenshot_bytes = await page.screenshot(path=path, full_page=full_page)
    page_size = await page.evaluate(
        """() => ({
            width: Math.max(document.documentElement.scrollWidth, document.body ? document.body.scrollWidth : 0),
            height: Math.max(document.documentElement.scrollHeight, document.body ? document.body.scrollHeight : 0)
        })"""
    )
    return {
        "path": path,
        "mime": "image/png",
        "bytes": len(screenshot_bytes),
        "full_page": full_page,
        "viewport": page.viewport_size,
        "page_size": page_size,
    }


_PLAYWRIGHT_SESSIONS: dict[str, dict[str, Any]] = {}
_PLAYWRIGHT_SESSIONS_LOCK = asyncio.Lock()


def _serializable_preview(value: Any, max_chars: int = DEFAULT_PLAYWRIGHT_MAX_CHARS) -> Any:
    if isinstance(value, str):
        clipped, truncated, remaining_chars = _truncate_text(value, max_chars)
        if not truncated:
            return clipped
        return {
            "content": clipped,
            "truncated": True,
            "remaining_chars": remaining_chars,
            "remaining_tokens": _count_remaining_tokens(value, max_chars),
        }
    try:
        json.dumps(value, ensure_ascii=False)
        return value
    except TypeError:
        return repr(value)


def _playwright_timeout_seconds(body: dict, default: int = DEFAULT_PLAYWRIGHT_SESSION_TIMEOUT) -> int:
    try:
        requested = int(body.get("timeout") or default)
    except (TypeError, ValueError):
        requested = default
    return max(1, min(requested, MAX_PLAYWRIGHT_SESSION_TIMEOUT))


def _resolve_file_list(value: Any) -> str | list[str]:
    if isinstance(value, list):
        return [_resolve_path(str(item)) for item in value]
    return _resolve_path(str(value))


async def _create_playwright_session(session_id: str, body: dict, timeout_ms: int) -> dict[str, Any]:
    from playwright.async_api import async_playwright

    headless = _coerce_bool(body.get("headless"), True)
    user_data_dir = (body.get("user_data_dir") or "").strip()
    storage_state_path = (body.get("storage_state_path") or "").strip()

    pw = await async_playwright().start()
    browser = None
    context = None
    try:
        if user_data_dir:
            context = await pw.chromium.launch_persistent_context(
                _resolve_path(user_data_dir),
                headless=headless,
            )
            browser = context.browser
        else:
            browser = await pw.chromium.launch(headless=headless)
            context_kwargs = {}
            if storage_state_path:
                resolved_state = _resolve_path(storage_state_path)
                if os.path.exists(resolved_state):
                    context_kwargs["storage_state"] = resolved_state
            context = await browser.new_context(**context_kwargs)

        page = context.pages[0] if context.pages else await context.new_page()
        page.set_default_timeout(timeout_ms)
        page.set_default_navigation_timeout(timeout_ms)

        return {
            "session_id": session_id,
            "playwright": pw,
            "browser": browser,
            "context": context,
            "page": page,
            "headless": headless,
            "user_data_dir": _resolve_path(user_data_dir) if user_data_dir else None,
            "storage_state_path": _resolve_path(storage_state_path) if storage_state_path else None,
            "persistent_context": bool(user_data_dir),
            "created_at": time.time(),
            "last_used_at": time.time(),
            "lock": asyncio.Lock(),
        }
    except Exception:
        with contextlib.suppress(Exception):
            if context is not None:
                await context.close()
        with contextlib.suppress(Exception):
            if browser is not None:
                await browser.close()
        with contextlib.suppress(Exception):
            await pw.stop()
        raise


async def _get_playwright_session(session_id: str, body: dict, timeout_ms: int) -> tuple[dict[str, Any], bool]:
    async with _PLAYWRIGHT_SESSIONS_LOCK:
        session = _PLAYWRIGHT_SESSIONS.get(session_id)
        if session is not None:
            session["last_used_at"] = time.time()
            return session, False

        session = await _create_playwright_session(session_id, body, timeout_ms)
        _PLAYWRIGHT_SESSIONS[session_id] = session
        return session, True


async def _close_playwright_session(session_id: str, save_storage_state: bool = True) -> dict:
    async with _PLAYWRIGHT_SESSIONS_LOCK:
        session = _PLAYWRIGHT_SESSIONS.pop(session_id, None)

    if session is None:
        return {"success": False, "error": f"session not found: {session_id}", "session_id": session_id}

    storage_state_path = session.get("storage_state_path")
    if save_storage_state and storage_state_path:
        with contextlib.suppress(Exception):
            await session["context"].storage_state(path=storage_state_path)

    errors = []
    try:
        await session["context"].close()
    except Exception as exc:
        errors.append(f"context: {exc}")

    if not session.get("persistent_context"):
        browser = session.get("browser")
        if browser is not None:
            try:
                await browser.close()
            except Exception as exc:
                errors.append(f"browser: {exc}")

    try:
        await session["playwright"].stop()
    except Exception as exc:
        errors.append(f"playwright: {exc}")

    return {
        "success": not errors,
        "session_id": session_id,
        "closed": True,
        "errors": errors,
    }


async def _wait_after_playwright_step(page, step: dict, timeout_ms: int) -> dict:
    waits: dict[str, Any] = {}

    wait_for = step.get("wait_for")
    if wait_for:
        await page.wait_for_selector(str(wait_for), timeout=timeout_ms)
        waits["selector"] = wait_for

    wait_for_text = step.get("wait_for_text")
    if wait_for_text:
        exact = _coerce_bool(step.get("exact"), False)
        await page.get_by_text(str(wait_for_text), exact=exact).wait_for(timeout=timeout_ms)
        waits["text"] = wait_for_text

    wait_for_url = step.get("wait_for_url")
    if wait_for_url:
        await page.wait_for_url(str(wait_for_url), timeout=timeout_ms)
        waits["url"] = wait_for_url

    wait_for_load_state = step.get("wait_for_load_state")
    if wait_for_load_state:
        await page.wait_for_load_state(str(wait_for_load_state), timeout=timeout_ms)
        waits["load_state"] = wait_for_load_state

    wait_ms = int(step.get("wait_ms") or 0)
    if wait_ms > 0:
        await page.wait_for_timeout(min(wait_ms, timeout_ms))
        waits["wait_ms"] = min(wait_ms, timeout_ms)

    return waits


async def _run_playwright_session_step(page, step: dict, body: dict, timeout_ms: int) -> dict:
    step_type = (step.get("type") or step.get("action") or "").strip().lower()
    if not step_type:
        return {"success": False, "error": "missing step type/action"}

    selector = step.get("selector") or body.get("selector") or ""
    value = step.get("value", body.get("value"))
    result: dict[str, Any] = {"success": True, "type": step_type}

    if step_type in {"goto", "navigate"}:
        url = (step.get("url") or body.get("url") or "").strip()
        if not url:
            return {"success": False, "type": step_type, "error": "url required"}
        response = await page.goto(url, timeout=timeout_ms)
        result.update({
            "url": page.url,
            "status": response.status if response else None,
        })

    elif step_type == "click":
        if not selector:
            return {"success": False, "type": step_type, "error": "selector required"}
        await page.click(selector, timeout=timeout_ms)
        result["selector"] = selector

    elif step_type == "fill":
        if not selector:
            return {"success": False, "type": step_type, "error": "selector required"}
        await page.fill(selector, "" if value is None else str(value), timeout=timeout_ms)
        result["selector"] = selector

    elif step_type == "check":
        if not selector:
            return {"success": False, "type": step_type, "error": "selector required"}
        await page.check(selector, timeout=timeout_ms)
        result["selector"] = selector

    elif step_type == "uncheck":
        if not selector:
            return {"success": False, "type": step_type, "error": "selector required"}
        await page.uncheck(selector, timeout=timeout_ms)
        result["selector"] = selector

    elif step_type in {"select", "select_option"}:
        if not selector:
            return {"success": False, "type": step_type, "error": "selector required"}
        options = step.get("values", value)
        selected = await page.select_option(selector, options, timeout=timeout_ms)
        result.update({"selector": selector, "selected": selected})

    elif step_type == "press":
        key = step.get("key") or value
        if not key:
            return {"success": False, "type": step_type, "error": "key/value required"}
        target = selector or "body"
        await page.press(target, str(key), timeout=timeout_ms)
        result.update({"selector": target, "key": key})

    elif step_type in {"upload", "set_input_files"}:
        if not selector:
            return {"success": False, "type": step_type, "error": "selector required"}
        files = step.get("file_paths") or step.get("files") or step.get("path") or value
        if not files:
            return {"success": False, "type": step_type, "error": "file path(s) required"}
        resolved_files = _resolve_file_list(files)
        await page.set_input_files(selector, resolved_files, timeout=timeout_ms)
        result.update({"selector": selector, "files": resolved_files})

    elif step_type == "wait_for_selector":
        if not selector:
            return {"success": False, "type": step_type, "error": "selector required"}
        state = step.get("state") or "visible"
        await page.wait_for_selector(selector, state=state, timeout=timeout_ms)
        result.update({"selector": selector, "state": state})

    elif step_type == "wait_for_text":
        text = step.get("text") or value
        if not text:
            return {"success": False, "type": step_type, "error": "text/value required"}
        exact = _coerce_bool(step.get("exact"), False)
        await page.get_by_text(str(text), exact=exact).wait_for(timeout=timeout_ms)
        result.update({"text": text, "exact": exact})

    elif step_type == "wait_for_url":
        url_pattern = step.get("url") or value
        if not url_pattern:
            return {"success": False, "type": step_type, "error": "url/value required"}
        await page.wait_for_url(str(url_pattern), timeout=timeout_ms)
        result["url_pattern"] = url_pattern

    elif step_type == "wait_for_load_state":
        state = step.get("state") or value or "load"
        await page.wait_for_load_state(str(state), timeout=timeout_ms)
        result["state"] = state

    elif step_type == "get_text":
        target = selector or "body"
        text = await page.inner_text(target, timeout=timeout_ms)
        result.update({
            "selector": target,
            "text": _serializable_preview(text, int(body.get("max_chars") or DEFAULT_PLAYWRIGHT_MAX_CHARS)),
        })

    elif step_type in {"snapshot", "inspect"}:
        result["snapshot"] = await _collect_playwright_snapshot(
            page=page,
            max_chars=int(body.get("max_chars") or DEFAULT_PLAYWRIGHT_MAX_CHARS),
            include_text=_coerce_bool(body.get("include_text"), True),
            include_elements=_coerce_bool(body.get("include_elements"), True),
            include_html=_coerce_bool(body.get("include_html"), False),
        )

    elif step_type == "screenshot":
        full_page = _coerce_bool(step.get("full_page", body.get("full_page")), True)
        result["screenshot"] = await _save_playwright_screenshot(page, {**body, **step}, full_page=full_page)

    elif step_type == "evaluate":
        script = step.get("script") or value
        if not script:
            return {"success": False, "type": step_type, "error": "script/value required"}
        evaluated = await page.evaluate(str(script), step.get("arg"))
        result["result"] = _serializable_preview(evaluated, int(body.get("max_chars") or DEFAULT_PLAYWRIGHT_MAX_CHARS))

    else:
        return {"success": False, "type": step_type, "error": f"unknown step type: {step_type}"}

    waits = await _wait_after_playwright_step(page, step, timeout_ms)
    if waits:
        result["waits"] = waits
    result["current_url"] = page.url
    return result


async def action_playwright_session(body: dict) -> dict:
    """
    Mantiene sesiones Playwright persistentes para flujos web largos.
    Requiere action, o actions para batch. session_id es opcional en start/goto/batch.
    """
    requested_action = (body.get("action") or ("batch" if body.get("actions") else "start")).strip().lower()
    timeout_s = _playwright_timeout_seconds(body)
    timeout_ms = timeout_s * 1000
    session_id = (body.get("session_id") or "").strip() or f"pw_{uuid.uuid4().hex[:12]}"

    if requested_action == "list":
        return {
            "success": True,
            "sessions": [
                {
                    "session_id": sid,
                    "url": data["page"].url,
                    "headless": data.get("headless"),
                    "created_at": data.get("created_at"),
                    "last_used_at": data.get("last_used_at"),
                    "user_data_dir": data.get("user_data_dir"),
                }
                for sid, data in _PLAYWRIGHT_SESSIONS.items()
            ],
        }

    if requested_action == "close":
        return await _close_playwright_session(
            session_id,
            save_storage_state=_coerce_bool(body.get("save_storage_state"), True),
        )

    try:
        session, created = await _get_playwright_session(session_id, body, timeout_ms)
        async with session["lock"]:
            page = session["page"]
            page.set_default_timeout(timeout_ms)
            page.set_default_navigation_timeout(timeout_ms)

            if requested_action == "batch":
                raw_steps = body.get("actions") or []
                if not isinstance(raw_steps, list) or not raw_steps:
                    return {"success": False, "error": "actions must be a non-empty list", "session_id": session_id}
                steps = raw_steps
            elif requested_action == "start":
                steps = [{"type": "goto", "url": body["url"]}] if body.get("url") else []
            else:
                steps = [{**body, "type": requested_action}]

            step_results = []
            for index, step in enumerate(steps, start=1):
                if not isinstance(step, dict):
                    step_results.append({"success": False, "index": index, "error": "step must be an object"})
                    break
                step_result = await _run_playwright_session_step(page, step, body, timeout_ms)
                step_result["index"] = index
                step_results.append(step_result)
                if not step_result.get("success"):
                    break

            storage_state_path = (body.get("storage_state_path") or session.get("storage_state_path") or "").strip()
            if storage_state_path and _coerce_bool(body.get("save_storage_state"), False):
                await session["context"].storage_state(path=_resolve_path(storage_state_path))

            include_snapshot = _coerce_bool(body.get("include_snapshot"), True)
            result = {
                "success": all(step.get("success") for step in step_results),
                "session_id": session_id,
                "created": created,
                "action": requested_action,
                "timeout": timeout_s,
                "url": page.url,
                "steps": step_results,
            }
            if include_snapshot:
                result["snapshot"] = await _collect_playwright_snapshot(
                    page=page,
                    max_chars=int(body.get("max_chars") or DEFAULT_PLAYWRIGHT_MAX_CHARS),
                    include_text=_coerce_bool(body.get("include_text"), True),
                    include_elements=_coerce_bool(body.get("include_elements"), True),
                    include_html=_coerce_bool(body.get("include_html"), False),
                )
            return result

    except Exception as e:
        return {"success": False, "error": str(e), "session_id": session_id, "action": requested_action}


async def action_playwright_navigate(body: dict) -> dict:
    """
    Navega y/o interactúa con una página web usando Playwright (Chromium headless).
    Requiere: url (str).
    Opcional:
      - action: 'navigate' (default) | 'snapshot' | 'inspect' | 'click' | 'fill' | 'screenshot' | 'get_text'
      - selector (str): selector CSS/XPath para click/fill/get_text
      - value (str): texto a introducir en fill
      - wait_for (str): selector a esperar antes de retornar
      - timeout (int, default 15, max 60): segundos
      - headless (bool, default True)
      - max_chars (int, default 12000): máximo texto/HTML devuelto
      - include_html/include_text/include_elements (bool): controla el snapshot
      - screenshot_mode: 'path' (default) | 'base64' | 'both'
    """
    from playwright.async_api import async_playwright

    url = (body.get("url") or "").strip()
    if not url:
        return {"success": False, "error": "missing required field: url"}

    action = (body.get("action") or "navigate").lower()
    selector = body.get("selector") or ""
    value = body.get("value") or ""
    wait_for = body.get("wait_for") or ""
    timeout_s = min(int(body.get("timeout") or 300), 600)
    timeout_ms = timeout_s * 1000
    headless = body.get("headless", True)
    max_chars = max(int(body.get("max_chars") or DEFAULT_PLAYWRIGHT_MAX_CHARS), 1)
    include_html = _coerce_bool(body.get("include_html"), False)
    include_text = _coerce_bool(body.get("include_text"), True)
    include_elements = _coerce_bool(body.get("include_elements"), True)
    screenshot_mode = (body.get("screenshot_mode") or "path").lower()
    if screenshot_mode not in {"path", "base64", "both"}:
        screenshot_mode = "path"

    try:
        async with async_playwright() as pw:
            browser = await pw.chromium.launch(headless=headless)
            page = await browser.new_page()
            await page.goto(url, timeout=timeout_ms)

            if wait_for:
                await page.wait_for_selector(wait_for, timeout=timeout_ms)

            result: dict = {"success": True, "url": page.url, "action": action}

            if action in {"navigate", "snapshot", "inspect"}:
                result["snapshot"] = await _collect_playwright_snapshot(
                    page=page,
                    max_chars=max_chars,
                    include_text=include_text,
                    include_elements=include_elements,
                    include_html=include_html,
                )

            elif action == "click":
                if not selector:
                    await browser.close()
                    return {"success": False, "error": "selector required for 'click'"}
                await page.click(selector, timeout=timeout_ms)
                result["snapshot"] = await _collect_playwright_snapshot(
                    page, max_chars, include_text, include_elements, include_html
                )

            elif action == "fill":
                if not selector:
                    await browser.close()
                    return {"success": False, "error": "selector required for 'fill'"}
                await page.fill(selector, value, timeout=timeout_ms)
                result["filled"] = value
                result["snapshot"] = await _collect_playwright_snapshot(
                    page, max_chars, include_text, include_elements, include_html
                )

            elif action == "get_text":
                if not selector:
                    await browser.close()
                    return {"success": False, "error": "selector required for 'get_text'"}
                text = await page.inner_text(selector, timeout=timeout_ms)
                result["text"], result["truncated"], result["remaining_chars"] = _truncate_text(text, max_chars)
                result["remaining_tokens"] = _count_remaining_tokens(text, max_chars)

            elif action == "screenshot":
                full_page = _coerce_bool(body.get("full_page"), True)
                screenshot = await _save_playwright_screenshot(page, body, full_page=full_page)
                result["screenshot"] = screenshot
                if screenshot_mode in {"base64", "both"}:
                    with open(screenshot["path"], "rb") as f:
                        result["screenshot_base64"] = base64.b64encode(f.read()).decode("ascii")
                result["snapshot"] = await _collect_playwright_snapshot(
                    page=page,
                    max_chars=max_chars,
                    include_text=include_text,
                    include_elements=False,
                    include_html=False,
                )

            else:
                await browser.close()
                return {"success": False, "error": f"unknown action: {action}"}

            await browser.close()
            return result

    except Exception as e:
        return {"success": False, "error": str(e), "url": url}


MAX_IMAGE_BYTES = 20 * 1024 * 1024
IMAGE_EXT_TO_MIME = {
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".gif": "image/gif", ".webp": "image/webp", ".bmp": "image/bmp",
}


def _resolve_image_source(source: str, explicit_mime: str | None) -> str:
    if source.startswith(("http://", "https://")) or source.startswith("data:image/"):
        return source

    candidate_path = _resolve_path(source)
    if os.path.isfile(candidate_path):
        size = os.path.getsize(candidate_path)
        if size > MAX_IMAGE_BYTES:
            raise ValueError(f"image exceeds {MAX_IMAGE_BYTES} bytes: {size}")
        ext = os.path.splitext(candidate_path)[1].lower()
        mime = explicit_mime or IMAGE_EXT_TO_MIME.get(ext) or mimetypes.guess_type(candidate_path)[0] or "image/png"
        with open(candidate_path, "rb") as f:
            raw = f.read()
        return f"data:{mime};base64,{base64.b64encode(raw).decode('ascii')}"

    cleaned = "".join(source.split())
    try:
        decoded = base64.b64decode(cleaned, validate=True)
    except Exception as exc:
        raise ValueError(f"source is not a valid path, URL or base64: {exc}") from None
    if len(decoded) > MAX_IMAGE_BYTES:
        raise ValueError(f"image exceeds {MAX_IMAGE_BYTES} bytes: {len(decoded)}")
    mime = explicit_mime or "image/png"
    return f"data:{mime};base64,{cleaned}"


def _extract_response_text(response) -> str:
    parts: list[str] = []
    for item in getattr(response, "output", []) or []:
        if getattr(item, "type", None) != "message":
            continue
        for chunk in getattr(item, "content", []) or []:
            if getattr(chunk, "type", None) == "output_text":
                text = getattr(chunk, "text", None)
                if isinstance(text, str) and text:
                    parts.append(text)
    return "\n".join(parts)


async def action_interpret_image(body: dict) -> dict:
    """
    Analiza una imagen con un modelo de visión vía Responses API.
    Requiere: source (str) — path local, URL http(s), data URI o base64 puro.
    Opcional: prompt (str), detail ('low'|'high'|'auto'), mime_type (str), model (str).
    """
    source = (body.get("source") or "").strip()
    if not source:
        return {"success": False, "error": "missing required field: source"}

    prompt = (body.get("prompt") or "").strip() or "Describe detalladamente el contenido de esta imagen."
    detail = body.get("detail", "auto")
    if detail not in ("low", "high", "auto"):
        detail = "auto"
    model = body.get("model") or "gpt-5.5"
    explicit_mime = (body.get("mime_type") or "").strip() or None

    try:
        image_url = await asyncio.to_thread(_resolve_image_source, source, explicit_mime)
    except Exception as e:
        return {"success": False, "error": f"cannot load image: {e}"}

    try:
        response = await openai_client.responses.create(
            model=model,
            input=[
                {
                    "role": "user",
                    "content": [
                        {"type": "input_text", "text": prompt},
                        {"type": "input_image", "image_url": image_url, "detail": detail},
                    ],
                }
            ],
        )
    except Exception as e:
        return {"success": False, "error": str(e), "model": model}

    return {
        "success": True,
        "model": model,
        "detail": detail,
        "interpretation": _extract_response_text(response),
    }


async def action_memory_query(body: dict):
    sql = (body.get("sql") or "").strip()
    if not sql:
        return {"success": False, "error": "Campo 'sql' requerido"}
    embed_text = (body.get("embed_text") or "").strip() or None
    return await RAG.execute_safe_query(sql, embed_text=embed_text)

def action_ask_user(body: dict):
    raise body["question"]

def action_save_preference(body: dict):
    def clean(value):
        if value is None:
            return None
        text = str(value).strip()
        return text or None

    def exact_or_none(value):
        if value is None:
            return None
        text = str(value)
        return None if text.strip() == "" else text

    action = clean(body.get("action"))
    old_text = exact_or_none(body.get("old_text"))
    new_text_raw = body.get("new_text")
    new_text = str(new_text_raw) if new_text_raw is not None else None
    preference = clean(body.get("preference"))
    replace_all = bool(body.get("replace_all", False))
    path = _resolve_path(user_preferences_path)

    # Backward compatibility for previous calls where agents sent empty
    # old_text/new_text plus preference to mean "add this preference".
    if action is None:
        if preference:
            action = "add"
        elif old_text and new_text is not None:
            action = "delete" if new_text == "" else "replace"

    try:
        os.makedirs(os.path.dirname(path), exist_ok=True)

        if action == "add":
            if not preference:
                preference = clean(new_text)
            if not preference:
                return {"success": False, "error": "preference is required for action='add'", "path": path}

            prefix = ""
            if os.path.exists(path) and os.path.getsize(path) > 0:
                prefix = "\n"
            with open(path, "a", encoding="utf-8") as f:
                f.write(f"{prefix} -{preference}")

            return {
                "success": True,
                "action": "add",
                "path": path,
                "file_hash": _file_hash(path),
            }

        if action not in {"replace", "delete"}:
            return {
                "success": False,
                "error": "action must be one of: add, replace, delete",
                "path": path,
            }

        if not old_text:
            return {"success": False, "error": f"old_text is required for action='{action}'", "path": path}

        if not os.path.exists(path):
            return {"success": False, "error": f"preferences file not found: {path}", "path": path}

        if action == "delete":
            replacement_text = ""
        else:
            if new_text is None:
                return {"success": False, "error": "new_text is required for action='replace'", "path": path}
            replacement_text = new_text

        with open(path, "r", encoding="utf-8", errors="replace") as f:
            content = f.read()

        count = content.count(old_text)
        if count == 0:
            return {
                "success": False,
                "error": "old_text not found in preferences",
                "path": path,
            }

        if not replace_all and count > 1:
            return {
                "success": False,
                "error": f"old_text found {count} times — set replace_all=true or provide more context to make it unique",
                "path": path,
                "occurrences": count,
            }

        new_content = (
            content.replace(old_text, replacement_text)
            if replace_all
            else content.replace(old_text, replacement_text, 1)
        )

        with open(path, "w", encoding="utf-8") as f:
            f.write(new_content)

        return {
            "success": True,
            "action": action,
            "path": path,
            "replacements": count if replace_all else 1,
            "file_hash": _file_hash(path),
        }
    except Exception as e:
        return {"success": False, "path": path, "error": str(e)}

# ── Agentes (proxy pass-through) ────────────────────────────────────

def _passthrough_agent(body: dict) -> str:
    return json.dumps(body, ensure_ascii=False)


# ── Dispatcher ───────────────────────────────────────────────────────

_AGENTS = [
    "ExecutorAgent", "WebSearchAgent", "DeviceManagerAgent",
    "PlaywrightSessionAgent", "MCPManagerAgent", "MemoryAgent",
]

ticket_dispatcher: dict[str, callable] = {
    **{name: _passthrough_agent for name in _AGENTS},
    "read_file": action_read_file,
    "file_hash": action_file_hash,
    "write_file": action_write_file,
    "edit_file": action_edit_file,
    "run_command": action_run_command,
    "run_python": action_run_python,
    "search_files": action_search_files,
    "web_fetch": action_web_fetch,
    "playwright_session": action_playwright_session,
    "playwright_navigate": action_playwright_navigate,
    "ask_user": action_ask_user,
    "save_preference": action_save_preference,
    "memory_query": action_memory_query,
    "interpret_image": action_interpret_image,
}
